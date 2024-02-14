import os
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Determine the directory of the script
script_dir = os.path.dirname(os.path.abspath(__file__))

# Path to the random sales data Excel file
input_file_path = os.path.join(script_dir, "random_sales_data.xlsx")

# Load the random sales data from the Excel file
df = pd.read_excel(input_file_path, index_col=0)

# Create the charts folder if it doesn't exist
charts_folder = os.path.join(script_dir, 'charts')
if not os.path.exists(charts_folder):
    os.mkdir(charts_folder)

# Create a new PowerPoint presentation
prs = Presentation()

# Create charts for each product
for product in df.columns:
    # Create a chart using the seaborn library
    sns.barplot(x=df.index, y=df[product])
    plt.title(product)
    plt.xlabel('Month')
    plt.ylabel('Sales')
    plt.xticks(rotation=45)
    plt.tight_layout()

    # Save the chart to the charts folder
    chart_file = f"{product}_chart.png"
    chart_path = os.path.join(charts_folder, chart_file)
    plt.savefig(chart_path)
    plt.close()

    # Add a slide to the PowerPoint presentation and insert the chart and title
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = product

    left = Inches(0.5)
    top = Inches(1)
    width = Inches(9)
    height = Inches(6)
    slide.shapes.add_picture(chart_path, left, top, width=width, height=height)

# Save the PowerPoint presentation in the same directory as the script
ppt_file = 'random_sales_data.pptx'
ppt_path = os.path.join(script_dir, ppt_file)
prs.save(ppt_path)

print(f"PowerPoint presentation saved to: {ppt_path}")
